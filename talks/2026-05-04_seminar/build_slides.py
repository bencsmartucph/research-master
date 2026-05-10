"""
Build PowerPoint slides for the 2026-05-04 seminar presentation.

Version 2 (2026-05-01): rebalanced for ~20 min @ 130 wpm, more visual,
honest framing of the damage cascade as theoretical synthesis rather
than empirically demonstrated chain. Speaker notes ~150 words/slide.

Run:
    python talks/2026-05-04_seminar/build_slides.py
"""

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(r'C:\Users\PKF715\Documents\claude_repos\Research_Master')
TALK_DIR = ROOT / 'talks' / '2026-05-04_seminar'
FIG_DIR = ROOT / 'outputs' / 'figures'

# Figure paths
FIG6 = FIG_DIR / 'fig6_cwed_country_slopes.png'
FIG7 = FIG_DIR / 'fig7_cwed_subcomponents.png'
FIG2 = FIG_DIR / 'fig2_rti_vs_antiimmig_by_regime.png'
FIG3 = FIG_DIR / 'fig3_marginal_effects.png'

# Colours
NAVY    = RGBColor(0x0B, 0x3D, 0x91)
RED     = RGBColor(0xC9, 0x30, 0x2C)
GREY    = RGBColor(0x6C, 0x75, 0x7D)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
PALE_BG = RGBColor(0xF6, 0xF8, 0xFA)
WARN_BG = RGBColor(0xFF, 0xF5, 0xE1)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, *, size=24, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, font='Calibri'):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tx


def add_multi_text(slide, x, y, w, h, parts, *, size=18, color=DARK,
                   align=PP_ALIGN.LEFT, font='Calibri'):
    """parts: list of (text, {'bold': bool, 'italic': bool, 'color': rgb, 'size': pt})"""
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, item in enumerate(parts):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        if i == 0:
            p = tf.paragraphs[0]
        elif text == '\n':
            p = tf.add_paragraph()
            continue
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get('size', size))
        run.font.bold = opts.get('bold', False)
        run.font.italic = opts.get('italic', False)
        run.font.color.rgb = opts.get('color', color)
        run.font.name = font
    return tx


def add_bullet_list(slide, x, y, w, h, items, *, size=18, color=DARK, font='Calibri'):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        level, text = (it if isinstance(it, tuple) else (0, it))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4); p.space_after = Pt(4)
        run = p.add_run()
        run.text = ('•  ' if level == 0 else '–  ') + text
        run.font.size = Pt(size - 2 * level)
        run.font.color.rgb = color
        run.font.name = font
    return tx


def add_title_bar(slide, text, *, color=NAVY):
    add_text(slide, 0.6, 0.35, 12, 0.7, text, size=28, bold=True, color=color)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.0),
                                   Inches(12.1), Inches(0.04))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = color


def add_image(slide, x, y, w, h, path):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        add_text(slide, x, y, w, 0.5, f"[FIGURE NOT FOUND: {path.name}]",
                 size=14, color=RED, italic=True)


def add_callout_box(slide, x, y, w, h, *, fill=PALE_BG, border_color=NAVY,
                    border_left_only=True, opacity=1.0):
    """Coloured rounded rectangle. Returns shape so you can add text on top."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if border_left_only:
        shape.line.fill.background()
        # Add left border via narrow rectangle
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x), Inches(y),
                                      Inches(0.06), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = border_color
        bar.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=NAVY, width=Pt(2)):
    """Draw an arrow from (x1,y1) to (x2,y2) in inches."""
    line = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = width
    # Add arrowhead via end_arrow
    line_format = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(line_format, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return line


def add_speaker_notes(slide, text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = text


# =====================================================================
# Slide 1 — Title
# =====================================================================
s = add_blank_slide()
add_text(s, 0, 2.6, 13.333, 1.2,
         "Dignity Is a Baseline",
         size=58, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0, 3.9, 13.333, 0.7,
         "Welfare Institutions and the Asymmetric Politics of Economic Disruption",
         size=22, italic=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.4, 13.333, 0.5,
         "Ben Smart",
         size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.95, 13.333, 0.4,
         "University of Copenhagen — Department of Economics",
         size=15, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, 0, 6.55, 13.333, 0.4,
         "Welfare State Seminar  ·  4 May 2026",
         size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)
add_speaker_notes(s, """[~20s]

Good afternoon. The thirty-second version: workers exposed to automation are
more anti-immigration than their material interests would predict. The gap is
bigger in some welfare states than others. The conventional explanation —
welfare buffers the backlash — is wrong about the mechanism. The mechanism is
asymmetric. I'll show you why.""")


# =====================================================================
# Slide 2 — The puzzle (2-column with callouts)
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "The puzzle")

# Left column — the empirical regularity
add_callout_box(s, 0.6, 1.4, 6.5, 2.3, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.85, 1.55, 6, 0.4, "Robust empirical regularity",
         size=15, bold=True, color=NAVY)
add_text(s, 0.85, 2.0, 6, 1.5,
         "Workers in routine-task-intensive (RTI) occupations across Europe disproportionately support populist radical right parties.",
         size=15, color=DARK)
add_text(s, 0.85, 3.2, 6, 0.5,
         "Gingrich (2019); Kurer (2020); Im et al. (2019); Autor et al. (2020)",
         size=11, italic=True, color=GREY)

add_callout_box(s, 0.6, 3.95, 6.5, 2.3, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.85, 4.1, 6, 0.4, "Cross-national variation",
         size=15, bold=True, color=NAVY)
add_text(s, 0.85, 4.55, 6, 1.5,
         "The same automation exposure converts into anti-immigration sentiment more readily in some welfare states than others.",
         size=15, color=DARK)
add_text(s, 0.85, 5.75, 6, 0.5,
         "Vlandas & Halikiopoulou (2022); Caselli et al. (2021)",
         size=11, italic=True, color=GREY)

# Right column — standard reading + my claim
add_callout_box(s, 7.5, 1.4, 5.3, 2.3, fill=WARN_BG, border_color=GREY)
add_text(s, 7.75, 1.55, 5, 0.4, "The standard reading",
         size=15, bold=True, color=GREY)
add_text(s, 7.75, 2.05, 5, 1.5,
         "Welfare buffers the backlash.",
         size=18, bold=True, color=DARK)
add_text(s, 7.75, 2.7, 5, 1,
         "Spend more, get less populism.",
         size=15, italic=True, color=GREY)

add_callout_box(s, 7.5, 3.95, 5.3, 2.8, fill=PALE_BG, border_color=NAVY)
add_text(s, 7.75, 4.1, 5, 0.4, "My claim",
         size=15, bold=True, color=NAVY)
add_text(s, 7.75, 4.55, 5, 1,
         "The standard reading is missing the mechanism.",
         size=16, bold=True, color=NAVY)
add_text(s, 7.75, 5.55, 5, 1.2,
         "What welfare communicates matters more than what welfare spends.",
         size=14, italic=True, color=DARK)

add_speaker_notes(s, """[~75s]

Two things are well documented. RTI workers vote radical right disproportionately.
And the pattern varies across welfare states.

The dominant explanation: welfare buffers the backlash. Quantity-based.
Compensation as cushion.

I'll argue today this account is missing the mechanism. Not wrong that welfare
matters; wrong about how. The dimension that matters is what welfare
communicates, not what it spends.""")


# =====================================================================
# Slide 3 — What buffering predicts + 3 pushbacks
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "What buffering predicts — and where it fails")

# Quote box
add_callout_box(s, 0.6, 1.3, 12.1, 1.0, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.9, 1.45, 11.5, 0.7,
         '"Following Ruggie\'s (1982) embedded liberalism bargain, generous compensation should dampen the political insecurity produced by economic openness."',
         size=15, italic=True, color=DARK)

add_text(s, 0.6, 2.5, 12, 0.4,
         "Operative variable: quantity. More spending → less populism. Symmetric mechanism.",
         size=15, color=DARK)

# Three failure cases — horizontal cards
y = 3.3
card_w = 4.0
card_h = 3.5
for i, (title, body, cite) in enumerate([
    ("Compensation that doesn't work",
     "Generous welfare states show no weaker automation→populism effect.",
     "Gingrich (2019)"),
    ("Compensation that backfires",
     "German coal phase-out: compensated communities had higher abstention and lower issue-owner support.",
     "Stutzmann (2025)"),
    ("Compensation that is resisted",
     "Workers refuse compensation that fully replaces income. Form & source matter independently of material content.",
     "Pelc (2025)"),
]):
    x = 0.6 + i * (card_w + 0.2)
    add_callout_box(s, x, y, card_w, card_h, fill=WARN_BG, border_color=RED)
    add_text(s, x + 0.2, y + 0.15, card_w - 0.4, 0.55, title,
             size=14, bold=True, color=RED)
    add_text(s, x + 0.2, y + 0.85, card_w - 0.4, 1.8, body,
             size=13, color=DARK)
    add_text(s, x + 0.2, y + card_h - 0.5, card_w - 0.4, 0.4,
             cite, size=12, italic=True, color=GREY)

add_speaker_notes(s, """[~75s]

Buffering rests on Ruggie. Compensation as quantity. Symmetric mechanism.

Three streams of evidence push back. Gingrich's cross-national: generous welfare
states don't show weaker automation-to-populism effects. Stutzmann on the German
coal phase-out: compensated communities had higher abstention. Pelc: workers
refuse compensation even when it fully replaces income.

Three different designs, same direction. Material compensation is doing less
of the political work than the buffering model assumes.""")


# =====================================================================
# Slide 4 — Recognition literature names it (3 quotes stacked)
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "What's missing — the recognition literature names it")

# Three stacked quote boxes
y = 1.3
quote_h = 1.7
for i, (q, attribution) in enumerate([
    ('"When relative societal decline rather than material hardship are at the heart of socially conservative resentment, traditional welfare policy may be an insufficient response..."',
     "Kurer (2020, p.1801)"),
    ('"This appeal to personal dignity is key to winning routine workers\' support. Perhaps even more than social protection, they demand economic AND cultural protection."',
     "Kurer & Palier (2019)"),
    ('Right-populist voters "care as much, or even more, about recognition as about redistribution."',
     "Gidron & Hall (2017, p.26)"),
]):
    yy = y + i * (quote_h + 0.2)
    add_callout_box(s, 0.6, yy, 12.1, quote_h, fill=PALE_BG, border_color=NAVY)
    add_text(s, 0.9, yy + 0.15, 11.5, 1.1, q,
             size=15, italic=True, color=DARK)
    add_text(s, 0.9, yy + quote_h - 0.45, 11.5, 0.4,
             "— " + attribution, size=13, bold=True, color=NAVY)

add_speaker_notes(s, """[~60s]

The recognition literature has been pointing at this for a decade.

Kurer: traditional welfare may be insufficient when relative decline drives
resentment. Kurer-Palier: the appeal is to dignity, not just protection. Gidron
and Hall: voters care as much about recognition as about redistribution.

The pattern: redistribution and recognition aren't fungible. The buffering
framework treats them as if they were. They aren't.""")


# =====================================================================
# Slide 5 — The argument (CENTRAL)
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "The argument")

add_text(s, 0.6, 1.35, 12, 0.6,
         "Welfare institutions are asymmetric in their political effects",
         size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Two-column callouts: failure / success
add_callout_box(s, 0.6, 2.2, 6.0, 3.0, fill=WARN_BG, border_color=RED)
add_text(s, 0.85, 2.35, 5.5, 0.4, "They CAN fail",
         size=16, bold=True, color=RED)
add_text(s, 0.85, 2.85, 5.5, 1.0,
         "By degrading recognition for vulnerable workers, conditioning three documented psychological mechanisms to engage in sequence:",
         size=12, color=DARK)
add_bullet_list(s, 1.0, 3.85, 5.3, 1.3, [
    "Identity switching  (Bonomi et al. 2021)",
    "Misattribution  (Gallego & Kurer 2022)",
    "Defensive othering  (Patrick 2016; Wagner 2022)",
], size=12)

add_callout_box(s, 6.85, 2.2, 6.0, 3.0, fill=PALE_BG, border_color=NAVY)
add_text(s, 7.1, 2.35, 5.5, 0.4, "They CANNOT symmetrically succeed",
         size=16, bold=True, color=NAVY)
add_text(s, 7.1, 2.85, 5.5, 2.0,
         "Producing solidarity requires constructive political work that welfare design alone cannot do.\n\nThe mirror-image mechanism does not exist.",
         size=13, color=DARK)

# Thesis box
thesis = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.6), Inches(5.5),
                             Inches(12.1), Inches(1.5))
thesis.fill.solid(); thesis.fill.fore_color.rgb = NAVY
thesis.line.fill.background()
ttf = thesis.text_frame; ttf.word_wrap = True
ttf.margin_left = Inches(0.4); ttf.margin_right = Inches(0.4)
ttf.margin_top = Inches(0.25)
tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
tr = tp.add_run()
tr.text = "Dignity is a baseline good."
tr.font.size = Pt(22); tr.font.bold = True
tr.font.color.rgb = WHITE; tr.font.name = 'Calibri'
tp2 = ttf.add_paragraph(); tp2.alignment = PP_ALIGN.CENTER
tr2 = tp2.add_run()
tr2.text = "Its absence damages. Its presence clears the ground for solidarity without producing it."
tr2.font.size = Pt(15); tr2.font.italic = True
tr2.font.color.rgb = WHITE; tr2.font.name = 'Calibri'

add_speaker_notes(s, """[~75s — central slide]

Welfare's political effects are asymmetric.

The literature has documented three psychological mechanisms separately. I argue
welfare design is the institutional condition that makes these mechanisms more
or less likely to engage. I'm not demonstrating the full chain empirically;
each step is independently documented, the institutional conditioning is my
theoretical claim.

The asymmetric piece: there's no mirror-image mechanism. Dignity-preserving
welfare clears space for solidarity but doesn't produce it. Solidarity requires
political construction welfare design alone can't supply.

That last sentence is the title in compressed form.""")


# =====================================================================
# Slide 6 — Documented vs argued (HONESTY SLIDE)
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "What's documented in the literature vs what I argue")

# Two columns
# LEFT — documented
add_callout_box(s, 0.6, 1.4, 6.0, 5.6, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.85, 1.55, 5.5, 0.45,
         "DOCUMENTED  (each step has its own literature)",
         size=13, bold=True, color=NAVY)
add_bullet_list(s, 0.85, 2.15, 5.5, 4.6, [
    "Stigmatising welfare implementation degrades class identity  (Soss 1999; Wagner 2022)",
    "Class identity loss → cultural identity activation  (Bonomi et al. 2021)",
    "Cultural identity → exclusionary attitudes  (Ballard-Rosa et al. 2022; Wu 2022)",
    "Conditional welfare → individual rather than structural attribution  (Alesina & Angeletos 2005)",
], size=13)

# RIGHT — argued
add_callout_box(s, 6.85, 1.4, 6.0, 5.6, fill=WARN_BG, border_color=RED)
add_text(s, 7.1, 1.55, 5.5, 0.45,
         "ARGUED  (my contribution, not yet proven)",
         size=13, bold=True, color=RED)
add_bullet_list(s, 7.1, 2.15, 5.5, 4.6, [
    "These three mechanisms operate as a connected sequence",
    "Welfare design is the institutional condition determining whether they engage",
    "The mechanism is asymmetric — no symmetric protective mechanism with comparable evidence",
    "The cross-national correlation is consistent with this; within-individual test belongs to the thesis follow-up",
], size=13)

add_speaker_notes(s, """[~60s — honesty slide, important]

I want to be careful here. Each step in the chain has its own literature. Soss
on welfare implementation. Bonomi et al on identity switching. Ballard-Rosa on
cultural identity to authoritarian values. Alesina-Angeletos on conditional
attribution.

What's mine is the synthesis: that these mechanisms are connected, that welfare
design is the institutional condition under which they engage together, that no
symmetric protective mechanism exists with comparable evidence.

The cross-national correlation is consistent with this argument. The within-
individual test belongs to the thesis. I'm not overclaiming.""")


# =====================================================================
# Slide 7 — Why no mirror image (3 reasons)
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "Why no symmetric mirror image — three reasons")

# Three columns
y = 1.4
card_h = 5.5
card_w = 4.0
reasons = [
    ("1.  Loss aversion",
     "Kahneman & Tversky (1979): losses loom larger than equivalent gains.",
     "A stigmatising encounter registers as a status loss. A dignity-preserving one registers as the ABSENCE of damage.",
     "Damage mobilises; absence of damage tends not to.",
     None),
    ("2.  Status is positional",
     "Recognition cannot be redistributed without losses to the currently-recognised.",
     "Dignity-preserving welfare REMOVES an obstacle to inclusive solidarity.",
     "It does not, by itself, CONSTRUCT the inclusion.",
     "Gidron & Hall (2017)"),
    ("3.  Investments are sticky",
     "Once a worker has built deservingness on critique of those below, the identity investment is costly to reverse.",
     "Pierson's (1994) positive feedback runs FORWARD into support.",
     "The cascade runs forward into harder-to-reverse opposition.",
     "Caveat: strongest theoretical claim, least directly tested."),
]
for i, (title, p1, p2, p3, footer) in enumerate(reasons):
    x = 0.6 + i * (card_w + 0.2)
    add_callout_box(s, x, y, card_w, card_h, fill=PALE_BG, border_color=NAVY)
    add_text(s, x + 0.2, y + 0.15, card_w - 0.4, 0.5, title,
             size=16, bold=True, color=NAVY)
    add_text(s, x + 0.2, y + 0.75, card_w - 0.4, 1.2, p1, size=12, color=DARK)
    add_text(s, x + 0.2, y + 1.95, card_w - 0.4, 1.2, p2, size=12, color=DARK)
    add_text(s, x + 0.2, y + 3.15, card_w - 0.4, 1.2, p3,
             size=12, color=DARK, bold=True)
    if footer:
        add_text(s, x + 0.2, y + card_h - 0.55, card_w - 0.4, 0.4, footer,
                 size=11, italic=True, color=GREY)

add_speaker_notes(s, """[~75s]

Three reasons the mechanism runs only one way.

First, loss aversion. Damage is psychologically heavier than equivalent
non-damage. Applies to dignity shocks just as material ones.

Second, status is positional. Recognition can't be redistributed the way
money can. Welfare clears space for inclusive solidarity but doesn't construct
it.

Third — and the one I'll flag honestly — defensive othering, once committed
to, is costly to reverse. Pierson's positive feedback runs forward into
support; the damage cascade runs forward into opposition that's harder to
reverse. This is the strongest theoretical claim and the least directly
tested. Worth flagging that.""")


# =====================================================================
# Slide 8 — Empirical setup
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "Empirical setup")

# Two-column
add_text(s, 0.6, 1.35, 6, 0.45, "Data", size=18, bold=True, color=NAVY)
add_bullet_list(s, 0.6, 1.85, 6, 2.2, [
    "European Social Survey rounds 6–9 (2012–2018)",
    "34 countries, N = 188,764",
    "15 Western European countries in CWED welfare-quality comparison",
], size=14)

add_text(s, 0.6, 4.2, 6, 0.45, "Identification", size=18, bold=True, color=NAVY)
add_bullet_list(s, 0.6, 4.7, 6, 2.5, [
    "Cross-level interactions: RTI × Welfare → attitudes",
    "Country-wave fixed effects, cluster-robust SEs",
    "Random-slope mixed models for cross-national heterogeneity",
    "Cross-sectional design — claim is consistency, not causation",
], size=14)

# Right side — two welfare measures contrasted
add_text(s, 7.0, 1.35, 6, 0.45, "Two competing welfare measures",
         size=18, bold=True, color=NAVY)

add_callout_box(s, 7.0, 2.0, 5.8, 1.6, fill=WARN_BG, border_color=GREY)
add_text(s, 7.25, 2.15, 5.4, 0.4, "ALMP spending (% GDP)",
         size=15, bold=True, color=GREY)
add_text(s, 7.25, 2.6, 5.4, 0.5,
         "What the buffering literature uses",
         size=13, italic=True, color=GREY)
add_text(s, 7.25, 3.1, 5.4, 0.5,
         "Captures welfare EFFORT.",
         size=13, color=DARK)

add_callout_box(s, 7.0, 3.85, 5.8, 1.8, fill=PALE_BG, border_color=NAVY)
add_text(s, 7.25, 4.0, 5.4, 0.4, "CWED decommodification",
         size=15, bold=True, color=NAVY)
add_text(s, 7.25, 4.45, 5.4, 0.5,
         "Esping-Andersen index",
         size=13, italic=True, color=GREY)
add_text(s, 7.25, 4.95, 5.4, 0.7,
         "The degree to which workers can sustain themselves WITHOUT dependence on the market.",
         size=13, color=DARK)

# Prediction box
add_text(s, 7.0, 6.0, 6, 1,
         "Asymmetric mechanism predicts:\nALMP fails; CWED predicts strongly.",
         size=14, italic=True, bold=True, color=NAVY)

add_speaker_notes(s, """[~60s]

ESS rounds 6 to 9, 34 countries, N=188,764. The headline analysis restricts
to the 15 Western European countries with CWED welfare-quality data.

Cross-level interactions, country-wave fixed effects, cluster-robust SEs.
Cross-sectional design. The claim is consistency with the asymmetric
mechanism, not causation.

Two competing welfare measures: ALMP spending — what the buffering literature
typically uses — and CWED decommodification, the Esping-Andersen index.

The asymmetric mechanism makes a sharp prediction: ALMP shouldn't predict the
slope; CWED should.""")


# =====================================================================
# Slide 9 — Cross-national pattern by regime
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "Cross-national pattern: 5 welfare regimes")

add_image(s, 1.0, 1.3, 8.5, 5.0, FIG2)

# Right panel — regime ordering
add_text(s, 9.8, 2.0, 3.2, 0.4, "Slope: RTI → anti-immig",
         size=13, bold=True, color=NAVY)
ty = 2.5
for regime, slope, is_top, is_bot in [
    ("Liberal",     "0.512", True, False),
    ("Southern",    "0.462", False, False),
    ("Continental", "0.443", False, False),
    ("Nordic",      "0.413", False, False),
    ("Eastern",     "0.263", False, True),
]:
    color = NAVY if is_top else (GREY if is_bot else DARK)
    bold = is_top
    add_text(s, 9.8, ty, 2, 0.35, regime,
             size=13, color=color, bold=bold)
    add_text(s, 11.5, ty, 1.5, 0.35, "β = " + slope,
             size=13, color=color, bold=bold)
    ty += 0.4

add_text(s, 0.6, 6.5, 12, 0.5,
         "Same RTI exposure, different conversion into exclusion. Liberal steepest, Nordic flattest.",
         size=14, italic=True, color=DARK, align=PP_ALIGN.CENTER)

add_speaker_notes(s, """[~45s]

Cross-national pattern by regime. Each panel a welfare regime; line is RTI
predicting anti-immigration. Liberal steepest at point five one two. Nordic
flattest at point four one three.

Same routine-task exposure, different conversion into exclusion. The buffering
reading would predict more spending → flatter slope. We test that next.""")


# =====================================================================
# Slide 10 — ALMP vs CWED (HEADLINE)
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "The empirical headline: ALMP vs CWED")

add_image(s, 2.5, 1.3, 8.3, 4.3, FIG6)

# Below — two contrast boxes
add_callout_box(s, 0.6, 5.85, 6.0, 1.4, fill=WARN_BG, border_color=GREY)
add_text(s, 0.85, 5.95, 5.5, 0.4, "ALMP spending",
         size=14, bold=True, color=GREY)
add_text(s, 0.85, 6.4, 5.5, 0.4,
         "r = +0.01,  p = 0.97  (n.s.)",
         size=15, bold=True, color=DARK)
add_text(s, 0.85, 6.85, 5.5, 0.35,
         "Spending effort is uncorrelated with the slope.",
         size=12, italic=True, color=DARK)

add_callout_box(s, 6.85, 5.85, 6.0, 1.4, fill=PALE_BG, border_color=NAVY)
add_text(s, 7.1, 5.95, 5.5, 0.4, "CWED decommodification",
         size=14, bold=True, color=NAVY)
add_text(s, 7.1, 6.4, 5.5, 0.4,
         "r = −0.85,  p < 0.001",
         size=15, bold=True, color=NAVY)
add_text(s, 7.1, 6.85, 5.5, 0.35,
         "Decommodification accounts for 72% of cross-national variation.",
         size=12, italic=True, color=DARK)

add_speaker_notes(s, """[~90s — empirical highlight]

This is the paper's most important empirical contrast. Same fifteen Western
European countries, two ways of measuring welfare.

ALMP spending: r equals plus point zero one. Essentially zero. Spending more
on labour market policies has nothing to do with the cross-national pattern.

CWED decommodification: r equals negative point eight five. Seventy-two
percent of the cross-national variation. UK lowest, Norway highest. Exactly
what the asymmetric mechanism predicts and what the buffering account cannot
explain.

The difference: ALMP captures effort. You can spend a lot on punitive
workfare. CWED captures decommodification — what the welfare state lets you
have, not what it costs to provide. Dignity travels along the second variable.""")


# =====================================================================
# Slide 11 — Sub-components decomposition (NEW)
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "Decomposing CWED: which dimension carries the signal?")

add_image(s, 0.5, 1.3, 7.3, 3.7, FIG7)

# Right — table
add_text(s, 8.0, 1.4, 5, 0.4,
         "Individual-level: RTI × component → anti-immigration",
         size=12, bold=True, color=NAVY)

# Mini table
table_data = [
    ("Component",   "β",       "p"),
    ("Unemployment","−0.053",  "<0.001 ★"),
    ("Sickness",    "−0.037",  "0.003"),
    ("Pensions",    "−0.019",  "0.066"),
    ("Composite",   "−0.051",  "<0.001"),
]
ty = 2.0
for r, row in enumerate(table_data):
    is_header = r == 0
    is_highlight = (r == 1)
    fcolor = NAVY if (is_header or is_highlight) else DARK
    fbold = is_header or is_highlight
    add_text(s, 8.0, ty, 1.7, 0.35, row[0], size=12, bold=fbold, color=fcolor)
    add_text(s, 9.7, ty, 1.0, 0.35, row[1], size=12, bold=fbold, color=fcolor)
    add_text(s, 10.7, ty, 1.7, 0.35, row[2], size=12, bold=fbold, color=fcolor)
    ty += 0.35

add_text(s, 8.0, 4.0, 5, 0.35,
         "N = 81,887; country-wave FE; cluster-robust SEs",
         size=10, italic=True, color=GREY)

# Theory check below
add_callout_box(s, 0.6, 5.4, 12.1, 1.7, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.85, 5.55, 11.5, 0.4,
         "Theory predicts UE > SK > PEN.    Observed at individual level: UE > SK > PEN.  ✓",
         size=14, bold=True, color=NAVY)
add_text(s, 0.85, 6.05, 11.5, 0.5,
         "The cascade fires through the point of economic vulnerability — where automation-exposed workers actually meet the welfare state.",
         size=13, color=DARK)
add_text(s, 0.85, 6.6, 11.5, 0.4,
         "Implication for thesis: Danish 2003/2006/2013 activation reforms should produce damage signatures; pension reforms should not.",
         size=12, italic=True, color=GREY)

add_speaker_notes(s, """[~75s — new finding]

Decomposing CWED into three sub-components: unemployment generosity, sickness,
pensions. Theory predicts unemployment should drive the result — that's where
automation-exposed routine workers meet the welfare state.

Individual-level: unemployment beta minus zero point zero five three, p less
than point zero zero one. Sickness intermediate. Pensions weakest, marginal.

Predicted UE greater than SK greater than PEN. Observed: same. The asymmetric
mechanism's specific prediction holds at the test that matters.

For thesis design: focus on unemployment benefit reforms — 2003, 2006, 2013
activation reforms. Pension reforms shouldn't show damage signatures.""")


# =====================================================================
# Slide 12 — Asymmetric confirmation
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "Same data, opposite outcome — the asymmetric confirmation")

# Comparison table
ty = 1.5
add_text(s, 0.6, ty, 6.5, 0.45, "Outcome", size=14, bold=True, color=GREY)
add_text(s, 7.5, ty, 3.5, 0.45, "RTI × Liberal", size=14, bold=True, color=GREY)
add_text(s, 11.0, ty, 2, 0.45, "p", size=14, bold=True, color=GREY)
ul = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.95),
                         Inches(12.1), Inches(0.03))
ul.line.fill.background(); ul.fill.solid(); ul.fill.fore_color.rgb = NAVY

ty = 2.15
add_text(s, 0.6, ty, 6.5, 0.5, "Anti-immigration",
         size=18, bold=True, color=NAVY)
add_text(s, 7.5, ty, 3.5, 0.5, "β = +0.127",
         size=18, bold=True, color=NAVY)
add_text(s, 11.0, ty, 2, 0.5, "0.003 ✓",
         size=18, bold=True, color=NAVY)

ty = 2.85
add_text(s, 0.6, ty, 6.5, 0.5, "Redistribution support",
         size=18, color=DARK)
add_text(s, 7.5, ty, 3.5, 0.5, "β = +0.011",
         size=18, color=DARK)
add_text(s, 11.0, ty, 2, 0.5, "0.285 (n.s.)",
         size=18, color=GREY)

ul2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.5),
                          Inches(12.1), Inches(0.03))
ul2.line.fill.background(); ul2.fill.solid(); ul2.fill.fore_color.rgb = NAVY

# Findings boxes
add_callout_box(s, 0.6, 4.0, 6.0, 2.8, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.85, 4.15, 5.5, 0.4, "The exclusion side is robust",
         size=14, bold=True, color=NAVY)
add_text(s, 0.85, 4.6, 5.5, 1.0,
         "Welfare context cleanly attenuates the conversion of vulnerability into exclusion.",
         size=13, color=DARK)
add_text(s, 0.85, 5.65, 5.5, 1.0,
         "Robust across every specification.",
         size=13, italic=True, color=DARK)

add_callout_box(s, 6.85, 4.0, 6.0, 2.8, fill=WARN_BG, border_color=RED)
add_text(s, 7.1, 4.15, 5.5, 0.4, "The solidarity side is null",
         size=14, bold=True, color=RED)
add_text(s, 7.1, 4.6, 5.5, 1.5,
         "Welfare context does NOT detectably moderate conversion into solidarity.",
         size=13, color=DARK)
add_text(s, 7.1, 5.7, 5.5, 1.0,
         "ISSP 2006 confirmation: same null on different sample/outcome/period (β=+0.010, p=0.55).",
         size=12, italic=True, color=DARK)

add_text(s, 0.6, 7.0, 12, 0.4,
         "This is what the asymmetric mechanism predicts.",
         size=15, italic=True, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_speaker_notes(s, """[~75s]

Same data, two outcomes. RTI predicts anti-immigration more strongly in
Liberal regimes than Nordic — beta point one two seven, p equals point zero
zero three. Robust across every specification.

RTI also predicts slightly higher redistribution support across all regimes,
but the cross-regime variation in that pathway is small, non-significant, and
in the wrong direction.

Welfare context attenuates exclusion. It does not detectably moderate
solidarity.

Two readings: measurement limitation, or substantive confirmation. The
supplementary ISSP test on different data with different outcome returns the
same null. I take the substantive reading. It's what the theory predicts.""")


# =====================================================================
# Slide 13 — Implications
# =====================================================================
s = add_blank_slide()
add_title_bar(s, "Implications")

implications = [
    (NAVY, "For welfare-state theory:",
     "the political consequences of welfare design travel along WHAT welfare communicates, not HOW MUCH it spends"),
    (NAVY, "For the cultural-vs-economic debate:",
     "cultural backlash isn't a rival explanation to economic disruption; it's what economic disruption LOOKS LIKE, cross-nationally, where welfare institutions are less decommodifying"),
    (NAVY, "For policy:",
     "dignity-preserving welfare is necessary but not sufficient for solidarity. Active solidarity requires political work welfare design alone can't do"),
    (RED, "For the thesis follow-up:",
     "Danish registry data on individuals before and after the 2003/2006/2013 activation reforms — testing within-individual whether conditionality shocks produce damage signatures"),
]
for i, (color, title, body) in enumerate(implications):
    y = 1.55 + i * 1.35
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y),
                              Inches(0.06), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(s, 0.85, y + 0.05, 12, 0.4, title,
             size=15, bold=True, color=color)
    add_text(s, 0.85, y + 0.45, 12, 0.7, body,
             size=14, color=DARK)

add_speaker_notes(s, """[~50s]

Four implications.

Welfare-state theory: dimension that matters is what welfare communicates.

Cultural-vs-economic: cultural backlash IS what economic disruption looks
like under thin decommodification.

Policy: dignity is necessary but not sufficient.

Thesis: Danish registry data, the 2003, 2006, 2013 activation reforms —
within-individual test of the cascade.""")


# =====================================================================
# Slide 14 — Closing
# =====================================================================
s = add_blank_slide()

add_text(s, 0, 1.8, 13.333, 1.2,
         "Dignity is a baseline",
         size=58, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s, 0, 3.5, 13.333, 0.55,
         "Its absence damages.",
         size=24, italic=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 0, 4.05, 13.333, 0.55,
         "Its presence clears the ground for solidarity.",
         size=24, italic=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 0, 4.6, 13.333, 0.55,
         "It does not, by itself, produce solidarity.",
         size=24, italic=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(s, 0, 6.0, 13.333, 0.55,
         "Thank you.",
         size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0, 6.65, 13.333, 0.4,
         "ben.smart@econ.ku.dk",
         size=13, color=GREY, align=PP_ALIGN.CENTER)

add_speaker_notes(s, """[~30s]

Dignity is a baseline. Its absence damages. Its presence clears the ground
for solidarity. It does not, by itself, produce solidarity.

The asymmetric mechanism is the technical version. The line is the moral
version. Both are saying the same thing.

Thank you. Happy to take questions.""")


# =====================================================================
# Backup slides — for Q&A
# =====================================================================

# Backup: marginal effects
s = add_blank_slide()
add_title_bar(s, "Backup: marginal effects")
add_image(s, 2.5, 1.3, 8.3, 4.5, FIG3)
add_text(s, 0.6, 6.0, 12, 0.5,
         "A 1-SD increase in RTI is associated with:",
         size=15, bold=True, color=DARK)
add_bullet_list(s, 0.6, 6.5, 12, 1.5, [
    "0.32 additional scale points of anti-immigration sentiment in Liberal regimes",
    "0.20 in Nordic regimes — gap is significant and substantively meaningful",
], size=14)

# Backup: Burgoon & Schakel
s = add_blank_slide()
add_title_bar(s, "Backup: Burgoon & Schakel (2022)")
add_text(s, 0.6, 1.3, 12, 0.6, "They find:", size=18, bold=True, color=NAVY)
add_text(s, 0.6, 1.85, 12, 0.7,
         "welfare generosity dampens anti-globalisation nationalism in European party platforms.",
         size=15, color=DARK)
add_callout_box(s, 0.6, 3.0, 12.1, 3.5, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.85, 3.15, 11.5, 0.5,
         "Apparent contradiction with my null on ALMP — resolved by unit of analysis:",
         size=15, bold=True, color=NAVY)
add_bullet_list(s, 0.85, 3.7, 11.5, 1.8, [
    "B&S measure platform language at the PARTY level",
    "I measure attitudinal slopes at the INDIVIDUAL level conditional on RTI",
    "Mechanisms differ: elite incentives + coalition arithmetic vs institutional encounter + self-concept",
], size=14)
add_text(s, 0.85, 5.5, 11.5, 0.9,
         "Both findings can be true. Welfare generosity may dampen the SUPPLY of anti-globalisation rhetoric while the DEMAND for exclusionary attitudes responds to a different dimension entirely.",
         size=13, italic=True, color=DARK)

# Backup: Denmark
s = add_blank_slide()
add_title_bar(s, "Backup: Denmark complication")
add_text(s, 0.6, 1.3, 12, 0.7,
         "Despite high CWED generosity, Denmark shows steeper RTI → exclusion slope (β=0.50) than Finland, Sweden, or Norway.",
         size=14, color=DARK)
add_callout_box(s, 0.6, 2.4, 12.1, 3.0, fill=PALE_BG, border_color=NAVY)
add_text(s, 0.85, 2.55, 11.5, 0.5,
         "Reading: confirmation, not anomaly.",
         size=18, bold=True, color=NAVY)
add_text(s, 0.85, 3.15, 11.5, 1.0,
         "Danish 'flexicurity' combines generous benefits with high labour market flexibility and active job search requirements — generous in transfers, demanding in activation.",
         size=13, italic=True, color=DARK)
add_text(s, 0.85, 4.3, 11.5, 1.0,
         "The asymmetric mechanism predicts conditionality and surveillance damage the self-concept EVEN WHEN transfer levels are high. Conditionality communicates suspicion.",
         size=13, color=DARK)
add_text(s, 0.6, 5.7, 12, 0.4,
         "Robustness: country-level finding survives all single-country exclusions.",
         size=12, italic=True, color=GREY)
add_text(s, 0.6, 6.1, 12, 0.4,
         "r = −0.717 even with both highest-leverage observations dropped (p = 0.006).",
         size=12, italic=True, color=GREY)

# Backup: limitations
s = add_blank_slide()
add_title_bar(s, "Backup: limitations")
add_bullet_list(s, 0.6, 1.4, 12.5, 5.5, [
    "Cross-sectional design cannot establish temporal ordering. Within-individual test belongs to the registry follow-up.",
    "Country-level confounders: Nordic high-CWED also have higher social trust, stronger unions, PR systems, lower ethnic heterogeneity.",
    "N = 15 country-level observations for the headline correlation. Individual-level Model 3 (β = −0.06, p = 0.015, N ≈ 82k) is the more defensible test.",
    "Loss aversion claim applies behavioural economics by analogy to dignity shocks. Supporting work consistent but doesn't test it directly.",
    "The damage cascade as a connected sequence is theoretical synthesis, not empirically demonstrated. Each step has independent evidence; the chain is my contribution.",
], size=14)


# =====================================================================
# Save
# =====================================================================
out = TALK_DIR / 'Dignity_Is_a_Baseline_2026-05-04_v2.pptx'
prs.save(str(out))
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}  (14 main + 4 backup)")
